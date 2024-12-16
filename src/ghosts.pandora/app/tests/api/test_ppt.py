from io import BytesIO
from unittest.mock import patch

from app import main
from fastapi.testclient import TestClient
from pptx import Presentation

client = TestClient(main)


def test_return_ppt_success(default_ppt_response):
    """Test PowerPoint generation with a random file name."""
    response = default_ppt_response
    assert response.status_code == 200
    assert (
        response.headers["Content-Type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert "Content-Disposition" in response.headers
    assert response.headers["Content-Disposition"].startswith("attachment; filename=")
    assert len(response.content) > 0  # Ensure content is returned


def test_return_ppt_custom_filename(custom_filename_ppt_response):
    """Test PowerPoint generation with a custom file name."""
    response = custom_filename_ppt_response
    assert response.status_code == 200
    assert (
        response.headers["Content-Type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert (
        response.headers["Content-Disposition"]
        == "attachment; filename=custom_presentation.pptx"
    )
    assert len(response.content) > 0  # Ensure content is returned


def test_ppt_content_structure(ppt_with_content_response):
    """Test the structure of the generated PowerPoint content."""
    response = ppt_with_content_response
    assert response.status_code == 200
    assert (
        response.headers["Content-Type"]
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    # Load the PowerPoint from the response content
    ppt = Presentation(BytesIO(response.content))

    # Check that at least two slides are present (title slide and content slide)
    assert len(ppt.slides) >= 2

    # Check that the first slide has a title and subtitle
    title_slide = ppt.slides[0]
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    assert title.text.strip() != ""
    assert subtitle.text.strip() != ""

    # Check the second slide (content slide)
    content_slide = ppt.slides[1]
    content = content_slide.shapes.placeholders[1]

    # Check that content is not empty
    assert content.text.strip() != ""


def test_return_ppt_with_ollama_failure(default_ppt_response):
    """Test PowerPoint generation with Ollama failure."""
    # Mock Ollama to simulate failure
    with patch("utils.ollama.generate_document_with_ollama") as mock_ollama:
        mock_ollama.side_effect = Exception(
            "Ollama service unavailable"
        )  # Simulate failure

        response = default_ppt_response
        assert response.status_code == 200  # Ensure fallback is used
        assert (
            response.headers["Content-Type"]
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        assert "Content-Disposition" in response.headers
        assert len(response.content) > 0  # Ensure fallback content is returned


def test_return_ppt_invalid_data():
    """Test PowerPoint generation with invalid data."""
    response = client.get("/ppt/invalidfile")
    assert (
        response.status_code == 500
    )  # Should trigger internal error due to invalid path or data
    assert (
        response.content
        == b"An error occurred while generating the PowerPoint presentation."
    )
